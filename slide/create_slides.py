import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color palette
    NAVY = RGBColor(15, 23, 42)      # #0f172a
    INDIGO = RGBColor(79, 70, 229)   # #4f46e5
    SLATE = RGBColor(71, 85, 105)    # #475569
    WHITE = RGBColor(255, 255, 255)  # #ffffff
    LIGHT_GRAY = RGBColor(248, 250, 252) # #f8fafc
    BORDER_COLOR = RGBColor(226, 232, 240) # #e2e8f0
    GREEN = RGBColor(16, 185, 129)   # #10b981
    RED = RGBColor(239, 68, 68)      # #ef4444
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # Helper: Set slide background
    # ----------------------------------------------------
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # ----------------------------------------------------
    # Helper: Add header title
    # ----------------------------------------------------
    def add_header(slide, title_text, category="PROJECT MANAGEMENT"):
        # Header shape accent
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.05)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = INDIGO
        shape.line.color.rgb = INDIGO
        
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.3))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = INDIGO
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, NAVY)
    
    # Big Title Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "AI-Powered Project Management Board"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Advanced Orchestration with LangGraph & Google Gemini API"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = INDIGO
    p2.alignment = PP_ALIGN.LEFT
    
    # Metadata Box
    meta_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.5))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    
    pm_para = tf_meta.paragraphs[0]
    pm_para.text = "Presenter: Project Development Team"
    pm_para.font.name = 'Segoe UI'
    pm_para.font.size = Pt(14)
    pm_para.font.color.rgb = RGBColor(148, 163, 184) # light slate
    
    tech_para = tf_meta.add_paragraph()
    tech_para.text = "Tech Stack: React 19 • FastAPI • LangGraph • SQLite • Telegram Bot API"
    tech_para.font.name = 'Segoe UI'
    tech_para.font.size = Pt(12)
    tech_para.font.italic = True
    tech_para.font.color.rgb = RGBColor(148, 163, 184)

    # ----------------------------------------------------
    # SLIDE 2: Project Introduction & Value Proposition
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_header(slide, "Introducing the AI Task Coordinator")
    
    # Columns Layout: Left (Overview), Right (Core Features Card)
    # Left Box
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Core Concept"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    p.space_after = Pt(12)
    
    p2 = tf_left.add_paragraph()
    p2.text = "Our platform transcends conventional Kanban tools. It acts as an autonomous coordinator, utilizing an AI Agent constructed via LangGraph state charts to automate management overhead, resolve schedule bottlenecks, and direct workloads seamlessly."
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(14)
    p2.font.color.rgb = SLATE
    p2.line_spacing = 1.3
    
    # Right Box (Decorative Container Card)
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.5)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = BORDER_COLOR
    
    right_content = slide.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.1), Inches(4.1))
    tf_right = right_content.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "High-Value Innovations"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(10)
    
    points = [
        ("✨ Auto-Allocation", "Analyzes developer workload & skill compatibility to distribute tasks."),
        ("⚠️ Cascading Timelines", "Shifts downstream dependent tasks automatically on delay reports."),
        ("🔄 Subtask Status Rollup", "Completes parent tasks instantly when all subtasks reach DONE."),
        ("🔒 Strict Role Controls", "Enforces distinct actions for PM, Developer, and QA views.")
    ]
    
    for title, desc in points:
        p_pt = tf_right.add_paragraph()
        p_pt.text = f"{title}: "
        p_pt.font.name = 'Segoe UI'
        p_pt.font.size = Pt(13)
        p_pt.font.bold = True
        p_pt.font.color.rgb = INDIGO
        p_pt.space_before = Pt(6)
        
        # Add description text on the same paragraph
        run = p_pt.add_run()
        run.text = desc
        run.font.name = 'Segoe UI'
        run.font.size = Pt(13)
        run.font.bold = False
        run.font.color.rgb = SLATE

    # ----------------------------------------------------
    # SLIDE 3: Technical Architecture & Stack
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_header(slide, "The Smart Tech Stack", "ARCHITECTURE")
    
    # 3 Column Layout for Stack
    cols = [
        ("FRONTEND INTERFACE", "React 19 & TypeScript", [
            "Vite compiler for rapid development & HMR",
            "Vanilla CSS for clean, lightweight, custom glassmorphic layout",
            "Lucide Icons for visual consistency",
            "Multi-View Dashboard (Kanban, Tree, Gantt, Calendar, Health, HR)"
        ]),
        ("BACKEND ENGINE", "FastAPI & Python", [
            "Asynchronous handlers using SQLAlchemy & aiosqlite",
            "FastAPI routers for modular REST endpoints",
            "APScheduler background runner for active deadline scanning",
            "Telegram Bot API integration for real-time notifications"
        ]),
        ("AI ORCHESTRATOR", "LangGraph & Gemini API", [
            "StateGraph architecture defining logical workflow loops",
            "Google Gemini API via LangChain integration",
            "Auto-Allocation & Cascading Shift decision engines",
            "Task-level Pair Programming chatbot assistant"
        ])
    ]
    
    for i, (title, subtitle, bullets) in enumerate(cols):
        left_pos = Inches(0.8 + (i * 3.95))
        
        # Background card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.0), Inches(3.75), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_COLOR
        
        content_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), Inches(3.35), Inches(4.1))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Segoe UI'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = INDIGO
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = 'Segoe UI'
        p_sub.font.size = Pt(14)
        p_sub.font.bold = True
        p_sub.font.color.rgb = NAVY
        p_sub.space_after = Pt(14)
        
        for bullet in bullets:
            p_bl = tf.add_paragraph()
            p_bl.text = "• "
            p_bl.font.name = 'Segoe UI'
            p_bl.font.size = Pt(11.5)
            p_bl.font.color.rgb = SLATE
            p_bl.space_after = Pt(6)
            
            run = p_bl.add_run()
            run.text = bullet
            run.font.name = 'Segoe UI'
            run.font.size = Pt(11.5)
            run.font.color.rgb = SLATE

    # ----------------------------------------------------
    # Helper: Fetch & Render Mermaid Image
    # ----------------------------------------------------
    import base64
    import urllib.request

    def add_mermaid_slide(slide, spec, title, category):
        add_header(slide, title, category)
        
        # Base64 encode the spec
        encoded = base64.urlsafe_b64encode(spec.encode("utf8")).decode("utf8")
        url = f"https://mermaid.ink/img/{encoded}"
        
        # Download image
        script_dir = os.path.dirname(os.path.abspath(__file__))
        temp_img = os.path.join(script_dir, f"temp_{title.lower().replace(' ', '_').replace('&', 'and')}.png")
        
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req) as res:
                with open(temp_img, 'wb') as f:
                    f.write(res.read())
            
            # Add image to slide
            # Widescreen is 13.333 x 7.5. Header takes top 2.0 inches.
            # Center horizontally: (13.333 - 10.0) / 2 = 1.666
            # Center vertically below header: (7.5 - 2.0 - 4.5) / 2 + 2.0 = 0.5 + 2.0 = 2.5
            slide.shapes.add_picture(temp_img, Inches(1.666), Inches(2.3), width=Inches(10.0))
        except Exception as e:
            print(f"Failed to render mermaid slide {title}: {e}")
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(3.0))
            p = tb.text_frame.paragraphs[0]
            p.text = f"Error rendering flowchart: {e}\n\nMermaid Code:\n{spec}"
            p.font.name = 'Segoe UI'
            p.font.size = Pt(12)
            p.font.color.rgb = RED

    # ----------------------------------------------------
    # SLIDE 4: System Architecture Flowchart
    # ----------------------------------------------------
    arch_spec = """graph TD
    classDef client fill:#ffffff,stroke:#4f46e5,stroke-width:2px,color:#0f172a;
    classDef server fill:#ffffff,stroke:#4f46e5,stroke-width:2px,color:#0f172a;
    classDef db fill:#ffffff,stroke:#475569,stroke-width:2px,color:#0f172a;
    classDef external fill:#ffffff,stroke:#10b981,stroke-width:2px,color:#0f172a;
    classDef alert fill:#ffffff,stroke:#ef4444,stroke-width:2px,color:#0f172a;

    A[React 19 Frontend Client]:::client
    B[FastAPI Python Server]:::server
    C[(SQLite Database via SQLAlchemy)]:::db
    D[Google Gemini API]:::external
    E[APScheduler Service]:::server
    F[Telegram Bot API]:::alert

    A <-->|HTTP REST / WebSockets| B
    B <-->|SQLAlchemy ORM| C
    B <-->|LangGraph State| D
    B <-->|Job Dispatcher| E
    E -->|Trigger Alerts| F
    C -.->|Read Status| E"""

    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_mermaid_slide(slide, arch_spec, "System Architecture & Data Flow", "ARCHITECTURE FLOWCHART")

    # ----------------------------------------------------
    # SLIDE 5: Supervisor Orchestrator Flowchart
    # ----------------------------------------------------
    supervisor_spec = """graph LR
    classDef input fill:#ffffff,stroke:#475569,stroke-width:2px,color:#0f172a;
    classDef supervisor fill:#ffffff,stroke:#4f46e5,stroke-width:2px,color:#0f172a;
    classDef agent fill:#ffffff,stroke:#10b981,stroke-width:2px,color:#0f172a;
    classDef output fill:#ffffff,stroke:#4f46e5,stroke-width:2px,color:#0f172a;

    1[1. User Request / Command]:::input
    2[2. Supervisor Router Node]:::supervisor
    3[3. DB Config Dynamic Prompts]:::input
    4[4. Selected Agent Execution]:::agent
    5[5. State Sync & Response UI]:::output

    1 --> 2
    2 <-->|Fetch Prompt Configurations| 3
    2 -->|Dynamic Dispatch / Tool Calls| 4
    4 -->|Update DB & Transition State| 5"""

    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_mermaid_slide(slide, supervisor_spec, "Supervisor Orchestrator Pipeline", "LANGGRAPH STATE FLOW")

    # ----------------------------------------------------
    # SLIDE 6: Sub-Agent Dispatch Flowchart
    # ----------------------------------------------------
    dispatch_spec = """graph TD
    classDef main fill:#ffffff,stroke:#0f172a,stroke-width:2px,color:#0f172a;
    classDef db fill:#ffffff,stroke:#475569,stroke-width:2px,color:#0f172a;
    classDef agent fill:#ffffff,stroke:#10b981,stroke-width:2px,color:#0f172a;
    classDef label fill:#f8fafc,stroke:#e2e8f0,stroke-width:1px,color:#475569;

    Request[Request from Frontend / Chat]:::main
    Supervisor[Supervisor Agent]:::main
    
    DB_Load[Load Prompt & Description from DB]:::label
    DB[(Agent Config DB)]:::db
    
    Route1[Analyze & Route]:::label
    Agent1[Sub-Agent: Decomposer]:::agent

    Route2[Analyze & Route]:::label
    Agent2[Sub-Agent: Allocator]:::agent

    Route3[Analyze & Route]:::label
    Agent3[Sub-Agent: Delay Shifter]:::agent

    Route4[Analyze & Route]:::label
    Agent4[Sub-Agent: Health Analyst]:::agent

    Route5[Analyze & Route]:::label
    Agent5[Sub-Agent: General Chat]:::agent

    Request --> Supervisor
    
    Supervisor --> DB_Load
    DB_Load --> DB
    
    Supervisor --> Route1
    Route1 --> Agent1
    
    Supervisor --> Route2
    Route2 --> Agent2
    
    Supervisor --> Route3
    Route3 --> Agent3
    
    Supervisor --> Route4
    Route4 --> Agent4
    
    Supervisor --> Route5
    Route5 --> Agent5"""

    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_mermaid_slide(slide, dispatch_spec, "Sub-Agent Coordination & Dispatch Flow", "ORCHESTRATOR ROUTING")

    # ----------------------------------------------------
    # SLIDE 7: Staff Roles & Even Task Distribution
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_header(slide, "Team Composition & Even Task Distribution", "TEAM RESOURCE PLANNING")
    
    # Introduce A, B, C
    card_width = Inches(3.75)
    card_height = Inches(4.5)
    
    members = [
        ("MEMBER A", "Project Manager (PM)", "Role: Overall Project Orchestrator", [
            "Responsible for project scheduling, defining tasks, and monitoring Gantt cascades.",
            "Assigned Tasks (Balanced workload):",
            "1. Project Planning & Req. Gathering (4 hrs)",
            "2. DB Schema Specification Review (2 hrs)",
            "3. Timeline Cascade & Progress Sign-off (3 hrs)"
        ], INDIGO),
        ("MEMBER B", "Frontend Developer", "Role: UI/UX Implementation", [
            "Responsible for building responsive interface components and layout styling.",
            "Assigned Tasks (Balanced workload):",
            "1. Dashboard Kanban View & D&D (6 hrs)",
            "2. Gantt Diagram SVG Layout (5 hrs)",
            "3. Interactive Forms & Chat Box (4 hrs)"
        ], GREEN),
        ("MEMBER C", "Backend Dev & QA", "Role: Core Logic & Quality Assurance", [
            "Responsible for FastAPI endpoints, DB engine, and validating task metrics.",
            "Assigned Tasks (Balanced workload):",
            "1. SQLAlchemy Model Setup & Migration (5 hrs)",
            "2. StateGraph Routing Logic (4 hrs)",
            "3. Comprehensive Testing & Validation (6 hrs)"
        ], RED)
    ]
    
    for i, (name, title, subtitle, points, theme_color) in enumerate(members):
        left_pos = Inches(0.8 + (i * 3.95))
        
        # Border box
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.0), card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_COLOR
        
        # Top banner accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(2.0), card_width, Inches(0.12))
        accent.fill.solid()
        accent.fill.fore_color.rgb = theme_color
        accent.line.fill.background()
        
        content_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), card_width - Inches(0.4), card_height - Inches(0.4))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = theme_color
        
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.name = 'Segoe UI'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = 'Segoe UI'
        p_sub.font.size = Pt(10)
        p_sub.font.italic = True
        p_sub.font.color.rgb = SLATE
        p_sub.space_after = Pt(10)
        
        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.font.name = 'Segoe UI'
            p_pt.font.size = Pt(10.5)
            p_pt.font.color.rgb = SLATE
            p_pt.space_after = Pt(5)
            
            if pt.startswith("Assigned Tasks") or pt.startswith("Role:"):
                p_pt.text = pt
                p_pt.font.bold = True
                p_pt.font.color.rgb = NAVY
            else:
                p_pt.text = "• " + pt
                p_pt.font.color.rgb = SLATE

    # ----------------------------------------------------
    # SLIDE 8: Demonstration Flow & Script Agenda
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, LIGHT_GRAY)
    add_header(slide, "Live Demo Scenarios Flow", "SYSTEM VERIFICATION")
    
    # Left column: Steps 1-3. Right column: Steps 4-5 + Telegram Alert Note
    # Left Box
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Part 1: Core Dashboard & Controls"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(10)
    
    steps_left = [
        ("1. Role Switcher & Strict Permissions", "Switch from Alice (PM) to Bob (Dev) to prove Bob cannot edit PM fields or approve Done tasks."),
        ("2. Real-time Kanban Drag & Drop", "Perform task status updates using drag and drop, and fill out the progress transition comment using AI suggestions."),
        ("3. Project & Task Auto-Decomposition", "Create a new project 'E-Commerce Mobile App' and trigger the AI Agent to generate Epics/Features/Tasks/Subtasks.")
    ]
    
    for num_title, desc in steps_left:
        p_step = tf_left.add_paragraph()
        p_step.text = num_title
        p_step.font.name = 'Segoe UI'
        p_step.font.size = Pt(13)
        p_step.font.bold = True
        p_step.font.color.rgb = INDIGO
        p_step.space_before = Pt(8)
        
        p_desc = tf_left.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Segoe UI'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = SLATE
        p_desc.space_after = Pt(6)
        
    # Right Box
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "Part 2: Intelligent Workflows"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(10)
    
    steps_right = [
        ("4. Auto Allocation & Gantt Cascade", "Use the ✨ Auto Allocate prompt. Report a delay on Task 35 to show the dependent tasks on the Gantt chart shift automatically."),
        ("5. QA Auto Routing & Task Blocking", "Mark all subtasks as Done to trigger auto-assignment to QA. Move a task to Blocked to trigger risk-mitigation banner alerts.")
    ]
    
    for num_title, desc in steps_right:
        p_step = tf_right.add_paragraph()
        p_step.text = num_title
        p_step.font.name = 'Segoe UI'
        p_step.font.size = Pt(13)
        p_step.font.bold = True
        p_step.font.color.rgb = INDIGO
        p_step.space_before = Pt(8)
        
        p_desc = tf_right.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Segoe UI'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = SLATE
        p_desc.space_after = Pt(6)
        
    # Telegram highlight card
    telegram_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(5.0), Inches(5.6), Inches(1.2))
    telegram_card.fill.solid()
    telegram_card.fill.fore_color.rgb = RGBColor(240, 249, 255) # light blue
    telegram_card.line.color.rgb = RGBColor(186, 230, 253)
    
    tg_box = slide.shapes.add_textbox(Inches(7.0), Inches(5.1), Inches(5.4), Inches(1.0))
    tf_tg = tg_box.text_frame
    tf_tg.word_wrap = True
    p_tg = tf_tg.paragraphs[0]
    p_tg.text = "🔔 Telegram Notifications Alert"
    p_tg.font.name = 'Segoe UI'
    p_tg.font.size = Pt(12)
    p_tg.font.bold = True
    p_tg.font.color.rgb = RGBColor(3, 105, 161) # dark blue
    
    p_tg_desc = tf_tg.add_paragraph()
    p_tg_desc.text = "Every status update, timeline shift, and approaching/overdue deadline automatically tags the corresponding user on Telegram for instant project visibility."
    p_tg_desc.font.name = 'Segoe UI'
    p_tg_desc.font.size = Pt(10)
    p_tg_desc.font.color.rgb = RGBColor(7, 89, 133)
 
    # ----------------------------------------------------
    # SLIDE 9: Summary & Thank You (Dark Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, NAVY)
    
    # Large Center Title
    center_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.333), Inches(3.5))
    tf = center_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Let's proceed to the Live Demonstration"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = INDIGO
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)
    
    # Save presentation
    script_dir = os.path.dirname(os.path.abspath(__file__))
    save_path = os.path.join(script_dir, "presentation.pptx")
    try:
        prs.save(save_path)
        print(f"Presentation saved successfully to: {os.path.abspath(save_path)}")
    except PermissionError:
        fallback_path = os.path.join(script_dir, "presentation_updated.pptx")
        prs.save(fallback_path)
        print(f"Permission denied on {save_path} (probably open in PowerPoint). Saved to: {os.path.abspath(fallback_path)}")
 
if __name__ == "__main__":
    create_presentation()
